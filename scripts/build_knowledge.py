#!/usr/bin/env python3
"""
Build data/cybersecurity.json from 8 source folders.
Logs all routing decisions to extraction-progress.log
"""
from __future__ import annotations

import json
import os
import re
import subprocess
import sys
import unicodedata
from html.parser import HTMLParser
from pathlib import Path
from xml.etree import ElementTree as ET

ROOT = Path(__file__).resolve().parents[1]
LOG_PATH = ROOT / "extraction-progress.log"
OUT_PATH = ROOT / "data" / "cybersecurity.json"
EXTRACT_DIR = ROOT / ".knowledge-build" / "extracted"

FOLDER_ORDER = [
    "COMS415", "CPRE430", "CPRE431", "CPRE489", "CPRE532", "CPRE536",
    "nmap-notes", "overthewire-solutions",
]

CLASS_META = {
    "COMS415": {
        "classId": "coms415", "classCode": "COMS 415",
        "className": "Software System Safety", "classType": "course",
    },
    "CPRE430": {
        "classId": "cpre430", "classCode": "CPRE 430",
        "className": "Networking Protocols & Security", "classType": "course",
    },
    "CPRE431": {
        "classId": "cpre431", "classCode": "CPRE 431",
        "className": "Information System Security", "classType": "course",
    },
    "CPRE489": {
        "classId": "cpre489", "classCode": "CPRE 489",
        "className": "Computer Network Design", "classType": "course",
    },
    "CPRE532": {
        "classId": "cpre532", "classCode": "CPRE 532",
        "className": "Information Warfare", "classType": "course",
    },
    "CPRE536": {
        "classId": "cpre536", "classCode": "CPRE 536",
        "className": "Computer and Network Forensics", "classType": "course",
    },
    "nmap-notes": {
        "classId": "nmap-notes", "classCode": "Nmap",
        "className": "Nmap Scanning Reference", "classType": "reference",
    },
    "overthewire-solutions": {
        "classId": "overthewire", "classCode": "Bandit",
        "className": "OverTheWire Bandit Solutions", "classType": "reference",
    },
}

SKIP_PATTERNS = [
    (lambda p: p.name.startswith("~$") and p.suffix.lower() in {".docx", ".odt"}, "Office temp/lock file"),
    (lambda p: p.suffix.lower() == ".zip", "zip archive"),
    (lambda p: p.suffix.lower() == ".mp4", "mp4 video"),
    (lambda p: p.suffix.lower() == ".o", "compiled object"),
    (lambda p: p.suffix.lower() == ".dat", "dat data file"),
    (lambda p: p.name == ".DS_Store", "macOS metadata"),
]

BINARY_NO_EXT = {"ruptimeclient", "ruptimeserver"}

log_lines: list[str] = []


def log(msg: str) -> None:
    log_lines.append(msg)
    print(msg)


def slugify(text: str, prefix: str = "") -> str:
    text = unicodedata.normalize("NFKD", text).encode("ascii", "ignore").decode()
    text = re.sub(r"[^a-zA-Z0-9]+", "-", text.lower()).strip("-")
    text = re.sub(r"-+", "-", text)[:60]
    return f"{prefix}-{text}" if prefix else text


def should_skip(path: Path) -> tuple[bool, str]:
    rel = path.name
    for pred, reason in SKIP_PATTERNS:
        if pred(path):
            return True, reason
    if not path.suffix and path.name.lower() in BINARY_NO_EXT:
        return True, "binary executable (no extension)"
    if not path.suffix and path.name != ".DS_Store":
        try:
            with open(path, "rb") as f:
                head = f.read(512)
            if b"\x00" in head[:128]:
                return True, "binary file (no extension)"
        except Exception:
            return True, "unreadable file"
    return False, ""


def setup_tesseract():
    try:
        import pytesseract
        for candidate in [
            "/opt/homebrew/bin/tesseract",
            "/usr/local/bin/tesseract",
            "/opt/homebrew/opt/tesseract/bin/tesseract",
        ]:
            if Path(candidate).exists():
                pytesseract.pytesseract.tesseract_cmd = candidate
                return True
        return bool(subprocess.run(["which", "tesseract"], capture_output=True).returncode == 0)
    except ImportError:
        return False


TESSERACT_OK = setup_tesseract()


def ocr_image(path: Path) -> tuple[str | None, str]:
    if not TESSERACT_OK:
        return None, "pytesseract/tesseract unavailable"
    try:
        import pytesseract
        from PIL import Image
        text = pytesseract.image_to_string(Image.open(path))
        text = re.sub(r"\n{3,}", "\n\n", text).strip()
        if len(text) < 20:
            return None, "OCR produced insufficient text"
        # garbled heuristic: >30% non-alphanumeric non-space
        weird = sum(1 for c in text if not c.isalnum() and c not in " \n\t|/-_.:$#@!%^&*()[]{}=+<>")
        if weird / max(len(text), 1) > 0.3:
            return None, "OCR output appears garbled"
        return text, "pytesseract"
    except Exception as e:
        return None, f"OCR failed: {e}"


def extract_pdf(path: Path) -> tuple[str | None, str]:
    try:
        import fitz
        doc = fitz.open(path)
        parts = []
        for page in doc:
            t = page.get_text().strip()
            if t:
                parts.append(t)
        text = "\n\n".join(parts).strip()
        if len(text) >= 50:
            return text, "PyMuPDF text layer"
        # try OCR per page
        if TESSERACT_OK:
            import pytesseract
            from PIL import Image
            ocr_parts = []
            for page in doc:
                pix = page.get_pixmap(dpi=150)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                ocr_parts.append(pytesseract.image_to_string(img))
            text = "\n\n".join(ocr_parts).strip()
            if len(text) >= 50:
                return text, "PyMuPDF+pytesseract OCR"
        return None, "PDF has no extractable text layer"
    except Exception as e:
        return None, f"PDF parse failed: {e}"


def extract_docx(path: Path) -> tuple[str | None, str]:
    try:
        from docx import Document
        doc = Document(path)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(c.text.strip() for c in row.cells if c.text.strip()))
        text = "\n".join(parts).strip()
        return (text, "python-docx") if text else (None, "empty docx")
    except Exception as e:
        return None, f"docx parse failed: {e}"


def extract_doc(path: Path) -> tuple[str | None, str]:
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(path)],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip(), "macOS textutil"
    except Exception:
        pass
    return None, "doc parse unavailable"


def extract_odt(path: Path) -> tuple[str | None, str]:
    try:
        result = subprocess.run(
            ["pandoc", "-f", "odt", "-t", "plain", str(path)],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip(), "pandoc"
    except Exception:
        pass
    try:
        ns = {"text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0"}
        tree = ET.parse(path)
        root = tree.getroot()
        parts = []
        for el in root.iter():
            if el.tag.endswith("}p") or el.tag.endswith("}h"):
                if el.text:
                    parts.append(el.text)
                for child in el:
                    if child.text:
                        parts.append(child.text)
                    if child.tail:
                        parts.append(child.tail)
        text = "\n".join(p.strip() for p in parts if p and p.strip())
        return (text, "odf XML parse") if text else (None, "empty odt")
    except Exception as e:
        return None, f"odt parse failed: {e}"


def extract_pptx(path: Path) -> tuple[str | None, str]:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                parts.append(f"Slide {i}:\n" + "\n".join(slide_text))
        text = "\n\n".join(parts).strip()
        return (text, "python-pptx") if text else (None, "empty pptx")
    except Exception as e:
        return None, f"pptx parse failed: {e}"


def extract_svg(path: Path) -> tuple[str | None, str]:
    try:
        raw = path.read_text(encoding="utf-8", errors="ignore")
        texts = re.findall(r"<text[^>]*>([^<]+)</text>", raw, re.I)
        texts += re.findall(r"<tspan[^>]*>([^<]+)</tspan>", raw, re.I)
        text = "\n".join(t.strip() for t in texts if t.strip())
        if text:
            return text, "SVG text elements"
        return None, "SVG text stored as vector paths (no text nodes)"
    except Exception as e:
        return None, f"SVG parse failed: {e}"


class _HTMLText(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts: list[str] = []
        self._skip = False

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style"):
            self._skip = True

    def handle_endtag(self, tag):
        if tag in ("script", "style"):
            self._skip = False
        if tag in ("p", "div", "section", "h1", "h2", "h3", "li", "br", "pre"):
            self.parts.append("\n")

    def handle_data(self, data):
        if not self._skip:
            self.parts.append(data)


def extract_html(path: Path) -> tuple[str | None, str]:
    try:
        raw = path.read_text(encoding="utf-8", errors="ignore")
        # preserve code blocks
        codes = re.findall(r"<pre[^>]*><code[^>]*>(.*?)</code></pre>", raw, re.S | re.I)
        parser = _HTMLText()
        parser.feed(raw)
        text = re.sub(r"\n{3,}", "\n\n", "".join(parser.parts))
        text = re.sub(r"[ \t]+", " ", text).strip()
        if codes:
            text += "\n\nCOMMANDS:\n" + "\n".join(
                re.sub(r"<[^>]+>", "", c).strip() for c in codes
            )
        return (text, "HTML tag strip") if text else (None, "empty html")
    except Exception as e:
        return None, f"html parse failed: {e}"


def extract_c_comments(path: Path) -> tuple[str | None, str]:
    try:
        raw = path.read_text(encoding="utf-8", errors="ignore")
        comments = re.findall(r"/\*.*?\*/", raw, re.S) + re.findall(r"//.*", raw)
        text = "\n".join(c.strip("/ *\n") for c in comments).strip()
        if len(text) > 80:
            return text, "C inline comments only"
        return None, "insufficient inline comments in C source"
    except Exception as e:
        return None, f"C read failed: {e}"


def extract_file(path: Path, folder: str) -> tuple[str | None, str, str]:
    """Returns (text, parser_used, skip_reason_if_none)"""
    skip, reason = should_skip(path)
    if skip:
        return None, "skip", reason

    ext = path.suffix.lower()
    rel = str(path.relative_to(ROOT / folder))

    if ext in {".md", ".txt"}:
        try:
            return path.read_text(encoding="utf-8", errors="replace").strip(), "direct read", ""
        except Exception as e:
            return None, "skip", f"read failed: {e}"
    if ext == ".html":
        return extract_html(path)[0], extract_html(path)[1], extract_html(path)[1] if not extract_html(path)[0] else ""
    if ext == ".pdf":
        t, m = extract_pdf(path)
        return t, m, m if not t else ""
    if ext == ".docx":
        t, m = extract_docx(path)
        return t, m, m if not t else ""
    if ext == ".doc":
        t, m = extract_doc(path)
        return t, m, m if not t else ""
    if ext == ".odt":
        t, m = extract_odt(path)
        return t, m, m if not t else ""
    if ext == ".pptx":
        t, m = extract_pptx(path)
        return t, m, m if not t else ""
    if ext == ".png":
        t, m = ocr_image(path)
        return t, m, m if not t else ""
    if ext == ".svg":
        t, m = extract_svg(path)
        return t, m, m if not t else ""
    if ext == ".xml":
        try:
            return path.read_text(encoding="utf-8", errors="replace").strip(), "xml read", ""
        except Exception as e:
            return None, "skip", str(e)
    if ext == ".py":
        try:
            return path.read_text(encoding="utf-8", errors="replace").strip(), "python read", ""
        except Exception as e:
            return None, "skip", str(e)
    if ext in {".c", ".h"}:
        t, m = extract_c_comments(path)
        return t, m, m if not t else ""

    return None, "skip", f"unrecognized format ({ext or 'no ext'})"


def first_sentence(text: str, max_len: int = 200) -> str:
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return ""
    m = re.match(r"^(.{20,}?[.!?])(?:\s|$)", text)
    if m and len(m.group(1)) <= max_len:
        return m.group(1).strip()
    return text[:max_len].rsplit(" ", 1)[0] + ("…" if len(text) > max_len else "")


def extract_commands_from_text(text: str) -> list[dict]:
    cmds: list[dict] = []
    seen: set[str] = set()

    # fenced code blocks
    for block in re.findall(r"```[\w]*\n(.*?)```", text, re.S):
        for line in block.strip().splitlines():
            line = line.strip()
            if line and not line.startswith("#") and len(line) < 200:
                if line not in seen:
                    seen.add(line)
                    cmds.append({"cmd": line, "explain": ""})

    # pre/code style lines after COMMANDS:
    if "COMMANDS:" in text:
        section = text.split("COMMANDS:", 1)[1]
        for line in section.splitlines():
            line = line.strip()
            if line and len(line) < 200 and line not in seen:
                seen.add(line)
                cmds.append({"cmd": line, "explain": ""})

    # backtick commands (shell-like)
    for m in re.finditer(r"`([^`]{2,120})`", text):
        cmd = m.group(1).strip()
        if re.search(r"^(nmap|ssh|ls|cat|cd|grep|find|man|whois|nslookup|fping|tracert|traceroute|tcpdump|wireshark|iptables|hydra|john|sqlmap|curl|wget|ping|chmod|sudo)\b", cmd, re.I):
            if cmd not in seen:
                seen.add(cmd)
                cmds.append({"cmd": cmd, "explain": ""})

    # table rows with | command |
    for line in text.splitlines():
        if "|" in line and re.search(r"\b(nmap|ssh|ls |cat |grep|find )\b", line, re.I):
            cells = [c.strip() for c in line.split("|") if c.strip()]
            for cell in cells:
                if re.match(r"^[a-zA-Z0-9_/. -]{3,80}$", cell) and " " in cell:
                    if cell not in seen:
                        seen.add(cell)
                        cmds.append({"cmd": cell, "explain": ""})

    return cmds[:12]


def infer_tags(title: str, text: str) -> list[str]:
    corpus = (title + " " + text).lower()
    tag_map = {
        "recon": ["recon", "footprint", "scan", "nmap", "enumeration"],
        "network": ["tcp", "udp", "ip", "routing", "dns", "protocol"],
        "cryptography": ["crypto", "encrypt", "hash", "tls", "certificate", "kerberos"],
        "exploitation": ["exploit", "overflow", "injection", "attack"],
        "forensics": ["forensic", "evidence", "memory", "disk"],
        "linux": ["linux", "bash", "ssh", "bandit", "shell"],
        "web": ["http", "web", "xss", "sql injection"],
        "wireless": ["wireless", "wifi", "802.11"],
        "ids": ["ids", "snort", "intrusion"],
        "firewall": ["firewall", "iptables", "acl"],
        "safety": ["safety", "hazard", "fault tree", "iec 61508"],
    }
    tags = [t for t, kws in tag_map.items() if any(k in corpus for k in kws)]
    return tags[:4] or ["general"]


def topics_from_markdown(text: str, source: str, class_id: str) -> list[dict]:
    topics = []
    module_m = re.search(r"^#\s+(.+)$", text, re.M)
    module_ctx = module_m.group(1).strip() if module_m else ""

    sections = re.split(r"\n(?=##\s+)", text)
    for sec in sections:
        if not sec.strip().startswith("##"):
            continue
        lines = sec.strip().splitlines()
        heading = re.sub(r"^##\s+", "", lines[0]).strip()
        body = "\n".join(lines[1:]).strip()
        if len(body) < 40:
            continue
        title = heading
        if module_ctx and module_ctx.lower() not in title.lower():
            short_mod = re.sub(r"^Module\s+\d+:\s*", "", module_ctx)
            title = f"{short_mod} — {heading}"
        topic_id = slugify(heading, class_id)
        cmds = extract_commands_from_text(body)
        # fill explain from following bullet if present
        for c in cmds:
            if not c["explain"]:
                c["explain"] = f"From {source}"
        topic = {
            "topicId": topic_id,
            "title": title[:100],
            "tags": infer_tags(title, body),
            "summary": first_sentence(body),
            "detail": body,
            "sourceFiles": [source],
        }
        if cmds:
            topic["commands"] = cmds
        topics.append(topic)
    return topics


def topics_from_otw_html(text: str, source: str, class_id: str, level_num: str) -> list[dict]:
    title_m = re.search(r"Level \d+:\s*([^\n]+)", text) or re.search(r"Hidden Files|SSH Login|", text)
    level_title = title_m.group(1).strip() if title_m and title_m.lastindex else source
    goal_m = re.search(r"Level Goal\s*(.+?)(?:Step-by-Step|Solution|$)", text, re.S | re.I)
    goal = goal_m.group(1).strip() if goal_m else ""
    cmds = extract_commands_from_text(text)
    for c in cmds:
        if not c["explain"]:
            c["explain"] = f"Bandit level {level_num} solution step"
    detail_parts = []
    if goal:
        detail_parts.append(f"**Goal:** {goal}")
    steps = re.split(r"\d+\.\s+", text)
    if len(steps) > 1:
        detail_parts.append("**Approach:**\n" + "\n".join(s.strip() for s in steps[1:4] if s.strip())[:2000])
    topic = {
        "topicId": slugify(f"level-{level_num}-{level_title}", class_id),
        "title": f"Bandit Level {level_num} — {level_title}",
        "tags": infer_tags(level_title, text) + ["linux", "ctf"],
        "summary": first_sentence(goal or text),
        "detail": "\n\n".join(detail_parts) or text[:2500],
        "sourceFiles": [source],
    }
    if cmds:
        topic["commands"] = cmds
    return [topic]


def topics_from_pdf_sections(text: str, source: str, class_id: str, doc_title: str = "") -> list[dict]:
    topics = []
    # split on numbered sections or ALL CAPS lines
    chunks = re.split(r"\n(?=\d+\.\s+[A-Z]|\n[A-Z][A-Z0-9 ,/&-]{8,}\n)", text)
    if len(chunks) <= 1:
        chunks = re.split(r"\n{2,}", text)
    for i, chunk in enumerate(chunks):
        chunk = chunk.strip()
        if len(chunk) < 80:
            continue
        lines = chunk.splitlines()
        title_line = lines[0].strip()
        if len(title_line) > 80:
            title_line = doc_title or source
        body = "\n".join(lines[1:]).strip() or chunk
        if len(body) < 60:
            continue
        topic_id = slugify(title_line[:50] or f"section-{i}", class_id)
        cmds = extract_commands_from_text(body)
        topic = {
            "topicId": topic_id,
            "title": (doc_title + " — " + title_line[:60]) if doc_title and doc_title not in title_line else title_line[:100],
            "tags": infer_tags(title_line, body),
            "summary": first_sentence(body),
            "detail": body[:4000],
            "sourceFiles": [source],
        }
        if cmds:
            topic["commands"] = cmds
        topics.append(topic)
        if len(topics) >= 8:
            break
    if not topics and len(text) > 100:
        topics.append({
            "topicId": slugify(source, class_id),
            "title": doc_title or source,
            "tags": infer_tags(doc_title, text),
            "summary": first_sentence(text),
            "detail": text[:4000],
            "sourceFiles": [source],
        })
    return topics


def topics_from_hex_packet(text: str, source: str, class_id: str) -> list[dict]:
    return [{
        "topicId": slugify("ethernet-ip-tcp-header-dissection", class_id),
        "title": "Wireshark Hex Dump — Ethernet/IP/TCP Header Fields",
        "tags": ["network", "protocol", "wireshark"],
        "summary": "Annotated hex capture mapping each field from Ethernet frame through IP options to TCP options — practice reading offsets without Wireshark's dissector pane.",
        "detail": "Field-by-field breakdown of a captured packet:\n\n" + text,
        "sourceFiles": [source],
    }]


def synthesize_from_extracted(folder: str, extracts: dict[str, str]) -> list[dict]:
    meta = CLASS_META[folder]
    class_id = meta["classId"]
    all_topics: list[dict] = []
    seen_ids: set[str] = set()

    def add_topics(new_topics: list[dict]) -> None:
        for t in new_topics:
            base_id = t["topicId"]
            tid = base_id
            n = 2
            while tid in seen_ids:
                tid = f"{base_id}-{n}"
                n += 1
            t["topicId"] = tid
            seen_ids.add(tid)
            all_topics.append(t)

    for rel, text in extracts.items():
        low = rel.lower()
        if low.endswith(".md"):
            add_topics(topics_from_markdown(text, rel, class_id))
        elif "overthewire" in folder and low.endswith(".html") and "level" in low:
            lvl = re.search(r"level(\d+)", low)
            add_topics(topics_from_otw_html(text, rel, class_id, lvl.group(1) if lvl else "0"))
        elif low == "hexpacket.txt":
            add_topics(topics_from_hex_packet(text, rel, class_id))
        elif low.endswith(".pdf"):
            doc_title = Path(rel).stem.replace("_", " ")
            add_topics(topics_from_pdf_sections(text, rel, class_id, doc_title))
        elif low.endswith((".docx", ".doc", ".pptx", ".odt")):
            doc_title = Path(rel).stem.replace("_", " ")
            add_topics(topics_from_pdf_sections(text, rel, class_id, doc_title))
        elif low.endswith(".txt") and "answer" in low:
            add_topics([{
                "topicId": slugify(rel, class_id),
                "title": Path(rel).stem,
                "tags": ["lab"],
                "summary": first_sentence(text),
                "detail": text,
                "sourceFiles": [rel],
            }])
        elif low.endswith(".png") and text:
            add_topics([{
                "topicId": slugify(Path(rel).stem, class_id),
                "title": f"Screenshot — {Path(rel).stem}",
                "tags": infer_tags(Path(rel).stem, text),
                "summary": first_sentence(text),
                "detail": text[:3000],
                "commands": extract_commands_from_text(text) or None,
                "sourceFiles": [rel],
            }])
        elif low.endswith(".py"):
            if len(text) > 100:
                add_topics([{
                    "topicId": slugify(rel, class_id),
                    "title": Path(rel).stem,
                    "tags": ["lab", "script"],
                    "summary": first_sentence(text),
                    "detail": f"```python\n{text[:3000]}\n```",
                    "sourceFiles": [rel],
                }])

    # dedupe by similar titles
    return all_topics


def process_folder(folder: str, data: dict) -> None:
    folder_path = ROOT / folder
    if not folder_path.is_dir():
        log(f"{folder}: ERROR — folder not found")
        return

    meta = CLASS_META[folder]
    extracts: dict[str, str] = {}
    skipped: list[str] = []
    used_files: list[str] = []
    parsers: dict[str, str] = {}

    for path in sorted(folder_path.rglob("*")):
        if not path.is_file():
            continue
        rel = str(path.relative_to(folder_path))
        text, parser, reason = extract_file(path, folder)

        if parser == "skip":
            skipped.append(f"{rel} ({reason})")
            log(f"  SKIP {rel}: {reason}")
            continue

        if text and len(text.strip()) >= 20:
            extracts[rel] = text
            used_files.append(rel)
            parsers[rel] = parser
            log(f"  EXTRACT {rel}: {parser} ({len(text)} chars)")
            EXTRACT_DIR.mkdir(parents=True, exist_ok=True)
            (EXTRACT_DIR / folder / rel).parent.mkdir(parents=True, exist_ok=True)
            (EXTRACT_DIR / folder / f"{rel}.txt").write_text(text, encoding="utf-8")
        else:
            skipped.append(f"{rel} ({reason or 'no content'})")
            log(f"  SKIP {rel}: {reason or 'no content'}")

    topics = synthesize_from_extracted(folder, extracts)

    # Update class name from README if present
    class_name = meta["className"]
    for rel, text in extracts.items():
        if rel.lower().endswith("readme.md"):
            m = re.search(r"^#\s+[^:]+:\s*(.+)$", text, re.M)
            if m:
                class_name = m.group(1).strip()

    class_entry = {
        **meta,
        "className": class_name,
        "topics": topics,
    }

    # Replace or append class
    existing = next((i for i, c in enumerate(data["classes"]) if c["classId"] == meta["classId"]), None)
    if existing is not None:
        data["classes"][existing] = class_entry
    else:
        data["classes"].append(class_entry)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    OUT_PATH.write_text(json.dumps(data, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")

    skip_summary = ", ".join(s.split(" (")[0] for s in skipped[:6])
    if len(skipped) > 6:
        skip_summary += f", +{len(skipped)-6} more"
    log(f"{folder}: {len(topics)} topics from {len(used_files)} files. Skipped: {skip_summary or 'none'}.")


def main():
    log("=== Knowledge extraction run ===")
    log(f"Tesseract available: {TESSERACT_OK}")
    data = {"category": "Cybersecurity", "classes": []}

    for folder in FOLDER_ORDER:
        log(f"\n--- Processing {folder} ---")
        try:
            process_folder(folder, data)
        except Exception as e:
            log(f"{folder}: ERROR — {e}")
            import traceback
            log(traceback.format_exc())

    LOG_PATH.write_text("\n".join(log_lines) + "\n", encoding="utf-8")
    log(f"\nWrote {OUT_PATH}")
    log(f"Log: {LOG_PATH}")


if __name__ == "__main__":
    main()
